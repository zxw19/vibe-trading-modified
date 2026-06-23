"""Universal document reader: dispatches by file extension.

Supported formats:
  - PDF (.pdf) — pypdfium2 + OCR fallback for image pages
  - Word (.docx) — python-docx (paragraphs + table cells)
  - Excel (.xlsx/.xls) — pandas preview, all sheets
  - PowerPoint (.pptx) — python-pptx (slide text)
  - Images (.png/.jpg/.jpeg/.gif/.bmp/.webp/.tiff) — OCR
  - Plain text (.txt/.md/.log/.json/.yaml/.yml/.toml/.ini/.cfg/.csv/.tsv/
                .html/.xml/.rst/.sql/.sh and common source-code extensions)

All handlers return the same JSON envelope: status, file, format, char_count,
truncated, text. PDF/Excel add format-specific metadata (pages, sheets, ...).
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Callable

from src.agent.progress import emit_progress
from src.agent.tools import BaseTool
from src.security.scanner import with_security_warnings
from src.tools.path_utils import safe_document_path

_MAX_CHARS = 15000
_MIN_TEXT_PER_PAGE = 50
_ENCODING_FALLBACK = ("utf-8", "utf-8-sig", "gbk", "gb2312", "big5", "latin-1")

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
_TEXT_EXTS = {
    # docs / structured
    ".txt", ".md", ".log", ".rst",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".csv", ".tsv", ".html", ".htm", ".xml",
    # source code (best-effort, LLM can read raw)
    ".py", ".pyi", ".js", ".jsx", ".ts", ".tsx",
    ".go", ".rs", ".java", ".kt", ".swift",
    ".c", ".h", ".cpp", ".hpp", ".cc",
    ".rb", ".php", ".pl", ".lua",
    ".sh", ".bash", ".zsh", ".ps1", ".bat",
    ".sql", ".r", ".m",
    ".dockerfile", ".makefile", ".cmake",
}

_ocr_engine = None


# ---------------- shared helpers ----------------

def _err(msg: str) -> str:
    return json.dumps({"status": "error", "error": msg}, ensure_ascii=False)


def _truncate(text: str) -> tuple[str, bool]:
    """Clip to _MAX_CHARS, return (text, was_truncated)."""
    if len(text) <= _MAX_CHARS:
        return text, False
    return text[:_MAX_CHARS] + f"\n\n... (truncated, total {len(text)} chars)", True


def _envelope(path: Path, fmt: str, text: str, **extra: Any) -> str:
    """Build the standard JSON response."""
    body, truncated = _truncate(text)
    payload: dict[str, Any] = {
        "status": "ok",
        "file": path.name,
        "format": fmt,
        "char_count": len(text),
        "truncated": truncated,
        "text": body,
    }
    payload.update(extra)
    payload = with_security_warnings(payload, fields=("text",))
    return json.dumps(payload, ensure_ascii=False)


def _get_ocr():
    """Lazily load RapidOCR. Raises ImportError if not installed."""
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR  # type: ignore
        _ocr_engine = RapidOCR()
    return _ocr_engine


def _ocr_image_array(img) -> str:
    """Run OCR on a numpy image; return joined lines or empty string."""
    try:
        ocr = _get_ocr()
    except ImportError:
        return ""
    result, _ = ocr(img)
    if not result:
        return ""
    return "\n".join(item[1] for item in result)


# ---------------- PDF ----------------

def _parse_pages(pages_str: str, total: int) -> list[int]:
    """Parse '1-10' / '5' / '1,3,5-8' into zero-based indices."""
    out: list[int] = []
    for part in pages_str.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start, end = part.split("-", 1)
            s = max(int(start.strip()) - 1, 0)
            e = min(int(end.strip()), total)
            out.extend(range(s, e))
        elif part.isdigit():
            out.append(int(part) - 1)
    return sorted(set(out))


def _read_pdf(path: Path, pages: str) -> str:
    """Extract PDF text; OCR pages with too little text."""
    try:
        import pypdfium2 as pdfium  # type: ignore
    except ImportError:
        return _err("pypdfium2 not installed; cannot read PDF")

    doc = pdfium.PdfDocument(str(path))
    try:
        total_pages = len(doc)
        targets = _parse_pages(pages, total_pages) if pages.strip() else list(range(total_pages))
        total_targets = len(targets)
        chunks: list[str] = []
        ocr_pages = 0
        for idx, i in enumerate(targets, start=1):
            if not 0 <= i < total_pages:
                continue
            page = doc[i]
            text = page.get_textpage().get_text_range().strip()
            if len(text) >= _MIN_TEXT_PER_PAGE:
                chunks.append(f"--- Page {i + 1} ---\n{text}")
                emit_progress(
                    "reading_pdf",
                    current=idx,
                    total=total_targets,
                    message=f"page {i + 1}/{total_pages}",
                )
                continue
            # OCR fallback for image pages
            bitmap = page.render(scale=300 / 72)
            img = bitmap.to_numpy()
            ocr_text = _ocr_image_array(img)
            if ocr_text.strip():
                chunks.append(f"--- Page {i + 1} [OCR] ---\n{ocr_text}")
                ocr_pages += 1
            elif text:
                chunks.append(f"--- Page {i + 1} ---\n{text}")
            emit_progress(
                "reading_pdf",
                current=idx,
                total=total_targets,
                message=f"page {i + 1}/{total_pages} (OCR)" if ocr_text.strip() else f"page {i + 1}/{total_pages}",
            )
        full = "\n\n".join(chunks)
        return _envelope(
            path, "pdf", full,
            total_pages=total_pages,
            pages_read=len(targets),
            ocr_pages=ocr_pages,
        )
    finally:
        doc.close()


# ---------------- DOCX ----------------

def _read_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        return _err("python-docx not installed; run: pip install python-docx")

    doc = docx.Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for t_idx, table in enumerate(doc.tables, start=1):
        parts.append(f"\n--- Table {t_idx} ---")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    return _envelope(path, "docx", "\n".join(parts), paragraphs=len(doc.paragraphs), tables=len(doc.tables))


# ---------------- Excel ----------------

def _read_excel(path: Path) -> str:
    try:
        import pandas as pd  # type: ignore
    except ImportError:
        return _err("pandas not installed; cannot read Excel")

    xls = pd.ExcelFile(path)
    parts: list[str] = []
    sheet_info: list[dict[str, Any]] = []
    total_sheets = len(xls.sheet_names)
    for idx, name in enumerate(xls.sheet_names, start=1):
        emit_progress(
            "reading_excel",
            current=idx,
            total=total_sheets,
            message=f"sheet {name}",
        )
        df = xls.parse(name, dtype=str)
        preview = df.head(100).to_string(index=False)
        parts.append(f"--- Sheet: {name} ({len(df)} rows × {len(df.columns)} cols) ---\n{preview}")
        sheet_info.append({"name": name, "rows": len(df), "cols": len(df.columns)})
    return _envelope(path, "excel", "\n\n".join(parts), sheets=sheet_info)


# ---------------- PPTX ----------------

def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return _err("python-pptx not installed; run: pip install python-pptx")

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return _envelope(path, "pptx", "\n".join(parts), slides=len(prs.slides))


# ---------------- Image OCR ----------------

def _read_image(path: Path) -> str:
    try:
        import numpy as np  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError:
        return _err("Pillow + numpy required for image OCR")

    try:
        img = np.array(Image.open(path).convert("RGB"))
    except Exception as exc:
        return _err(f"Failed to open image: {exc}")

    text = _ocr_image_array(img)
    if not text.strip():
        return _envelope(path, "image", "", note="OCR returned no text (engine missing or empty image)")
    return _envelope(path, "image", text)


# ---------------- Plain text ----------------

def _read_text(path: Path) -> str:
    """Read a text-like file with encoding fallback."""
    data = path.read_bytes()
    last_err: Exception | None = None
    for enc in _ENCODING_FALLBACK:
        try:
            decoded = data.decode(enc)
            return _envelope(path, "text", decoded, encoding=enc, size=len(data))
        except UnicodeDecodeError as exc:
            last_err = exc
    return _err(f"Failed to decode file with any of {_ENCODING_FALLBACK}: {last_err}")


# ---------------- Dispatcher ----------------

_HANDLERS: dict[str, Callable[[Path], str]] = {
    ".docx": _read_docx,
    ".xlsx": _read_excel,
    ".xls": _read_excel,
    ".pptx": _read_pptx,
}


def read_document(file_path: str, pages: str = "") -> str:
    """Read any supported document; dispatch by extension.

    Args:
        file_path: Absolute path to the file.
        pages: Only used for PDF — e.g. "1-10", "5", "1,3,5-8"; empty = all.

    Returns:
        JSON envelope: status, file, format, char_count, truncated, text,
        plus format-specific metadata (total_pages, sheets, etc.).
    """
    try:
        path = safe_document_path(file_path)
    except ValueError as exc:
        return _err(str(exc))
    if not path.exists():
        return _err(f"File not found: {file_path}")
    if not path.is_file():
        return _err(f"Not a file: {file_path}")

    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _read_pdf(path, pages)
        if ext in _HANDLERS:
            return _HANDLERS[ext](path)
        if ext in _IMAGE_EXTS:
            return _read_image(path)
        if ext in _TEXT_EXTS or ext == "":
            return _read_text(path)
        # Unknown extension: best-effort text read
        return _read_text(path)
    except Exception as exc:
        return _err(f"{type(exc).__name__}: {exc}")


class DocReaderTool(BaseTool):
    """Universal document reader — PDF/Word/Excel/PowerPoint/images/text."""

    name = "read_document"
    description = (
        "Read a document of any common format: PDF, Word (.docx), Excel "
        "(.xlsx/.xls), PowerPoint (.pptx), images (OCR), or plain text "
        "(txt/md/json/yaml/csv/html/code files). Returns extracted text in "
        "a unified JSON envelope. For PDFs, accepts an optional `pages` range."
    )
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Absolute path to the file."},
            "pages": {
                "type": "string",
                "description": "PDF only: page range (e.g. '1-10', '5', '1,3,5-8'). Ignored for other formats.",
                "default": "",
            },
        },
        "required": ["file_path"],
    }
    repeatable = True

    def execute(self, **kwargs: Any) -> str:
        return read_document(kwargs["file_path"], kwargs.get("pages", ""))
